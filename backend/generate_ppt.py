from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "UX Guardian"
    subtitle.text = "The Agentic Platform for Conversational UX & Accessibility"
    
    # Custom styling
    title.text_frame.paragraphs[0].font.size = Pt(60)
    title.text_frame.paragraphs[0].font.bold = True
    subtitle.text_frame.paragraphs[0].font.size = Pt(28)

    slides_data = [
        {
            "title": "The Problem Statement",
            "content": [
                "Static Dashboards Are Dead (Conversational UX)",
                "Traditional tools output massive, unreadable spreadsheets.",
                "Developers are overwhelmed by jargon; executives can't see financial impact.",
                "Accessibility is treated as an afterthought because it's hard to visualize.",
                "Users shouldn't have to navigate 5 menus. Conversation IS the interface."
            ]
        },
        {
            "title": "Our Solution & Features",
            "content": [
                "Automated Multi-Modal Auditing",
                "Vision AI scans both code (DOM) and visual hierarchy (Screenshot).",
                "Generative UI Command Center",
                "Press Cmd+K to bypass dashboards. Chat with AI to render React widgets dynamically."
            ]
        },
        {
            "title": "Proving Accessibility (A11y)",
            "content": [
                "Empathy-Driven Development",
                "The Empathy Simulator: Forces navigation via a synthetic screen-reader.",
                "The Voice Assistant: Conversational UI allowing hands-free interaction.",
                "Powered by Smallest.ai AWAAZ: Ultra-low latency, hyper-realistic TTS voices.",
                "We don't just audit for accessibility; our platform practices it natively."
            ]
        },
        {
            "title": "The Technical Stack",
            "content": [
                "Modern, Asynchronous Micro-Architecture",
                "Frontend: React 18, TypeScript, Tailwind CSS, Zustand, Framer Motion",
                "Backend: Python, FastAPI, Playwright (Headless Browser Automation)",
                "AI/LLM: Google Gemini 1.5 Pro & Flash, Pydantic (JSON schema enforcement)",
                "Infrastructure: Docker & Docker Compose"
            ]
        },
        {
            "title": "Technical Feasibility",
            "content": [
                "Built for Speed and Reliability",
                "Low Latency: Gemini 1.5 Flash drops multi-modal audit times to <10 seconds.",
                "Zero-Cost Edge Compute: Voice Assistant & Empathy Simulator run entirely client-side (0ms latency).",
                "Anti-Hallucination: Strict Pydantic schemas ensure AI returns parseable JSON arrays."
            ]
        },
        {
            "title": "Business Impact",
            "content": [
                "Calculating Revenue at Risk",
                "The ROI of UX: Algorithm correlates heuristic violations to estimated Conversion Drop-off.",
                "Actionable Metrics: Translates bad UX into hard dollars lost based on Traffic/AOV.",
                "Industry Benchmarking: Compare scores dynamically against E-commerce, SaaS, or Media averages."
            ]
        },
        {
            "title": "Conclusion & Demo",
            "content": [
                "Thank You",
                "Ready to experience the future of UX Auditing?",
                "We didn't build a dashboard. We built an Agentic Platform.",
                "Switching directly to the Live Demo..."
            ]
        }
    ]

    for data in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1]) # Title and Content
        title = slide.shapes.title
        title.text = data["title"]
        title.text_frame.paragraphs[0].font.size = Pt(40)
        title.text_frame.paragraphs[0].font.bold = True
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        
        for i, point in enumerate(data["content"]):
            if i == 0:
                p = tf.paragraphs[0]
                p.text = point
            else:
                p = tf.add_paragraph()
                p.text = point
                p.level = 1
            
            p.font.size = Pt(24) if p.level == 0 else Pt(20)

    prs.save('UX_Guardian_Pitch.pptx')

if __name__ == "__main__":
    create_presentation()
